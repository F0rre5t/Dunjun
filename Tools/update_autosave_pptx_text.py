"""Update text only inside the autosaved PPT — keep its existing template/shapes."""
from pathlib import Path
import shutil
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

FOLDER = Path(r"C:\Users\you\Rouge\Presentation")
BG = RGBColor(0x14, 0x18, 0x1C)
ACCENT = RGBColor(0xD4, 0xA3, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xF2, 0xF0, 0xEA)
MUTED = RGBColor(0xA8, 0xB0, 0xB8)
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def find_autosave() -> Path:
    # Exact WPS/Office autosave name pattern; ignore helper exports
    matches = []
    for p in FOLDER.iterdir():
        if p.name.startswith("~$"):
            continue
        if p.suffix.lower() != ".pptx":
            continue
        if "updated" in p.name.lower() or "AUTOSAVE_updated" in p.name:
            continue
        if "Dynamic_Difficulty_Video_Slides" in p.name and "自动保存" in p.name:
            matches.append(p)
    if not matches:
        raise FileNotFoundError("Autosaved PPTX not found (…[自动保存的].pptx)")
    return matches[0]


def set_shape_text(shape, text: str, *, size=None, bold=None, color=None):
    tf = shape.text_frame
    first = None
    for p in tf.paragraphs:
        for r in p.runs:
            first = r
            break
        if first:
            break

    default_size = size or (first.font.size if first and first.font.size else Pt(20))
    default_bold = bold if bold is not None else (bool(first.font.bold) if first else False)
    default_name = first.font.name if first and first.font.name else "Calibri"
    if color is None:
        try:
            default_color = first.font.color.rgb if first and first.font.color.type else TEXT
        except Exception:
            default_color = TEXT
    else:
        default_color = color

    lines = text.split("\n")
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag == qn("a:p"):
            txBody.remove(child)

    for line in lines:
        p_elem = parse_xml(
            '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        txBody.append(p_elem)

    for i, line in enumerate(lines):
        p = tf.paragraphs[i]
        run = p.add_run()
        run.text = line
        run.font.name = default_name
        run.font.size = default_size
        run.font.bold = default_bold
        run.font.color.rgb = default_color
        p.space_after = Pt(8)


def textboxes(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]


def insert_game_slide(prs, after_index: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(91440), prs.slide_height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.8))
    r = title.text_frame.paragraphs[0].add_run()
    r.text = "The Game"
    r.font.name = "Calibri"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    lines = [
        "PCG",
        "• Random-walk rooms · walls / doors from neighbours",
        "• Start / shop / boss by rules · enemies by depth",
        "• Spike chance & shape scale with progress",
        "",
        "Progression",
        "• Relics + gold shops → different builds each run",
        "• Route / build / resources vary → state swings hard",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.bold = line in ("PCG", "Progression")
        run.font.color.rgb = WHITE if line in ("PCG", "Progression") else TEXT
        p.space_after = Pt(6)

    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst)
    new_entry = entries[-1]
    sldIdLst.remove(new_entry)
    ref = list(sldIdLst)[after_index]
    ref.addnext(new_entry)
    return slide


def main():
    path = find_autosave()
    print(f"Opening: {path.name}")
    prs = Presentation(str(path))
    print(f"Original slides: {len(prs.slides)}")

    # If this file was already partially updated (13 slides), refuse and ask for clean source
    # We expect the original autosave with 12 slides.
    if len(prs.slides) != 12:
        raise RuntimeError(
            f"Expected 12-slide autosave template, found {len(prs.slides)}. "
            "Restore/close the original [自动保存的] file and retry."
        )

    # 1 Title
    s = prs.slides[0]
    tbs = textboxes(s)
    set_shape_text(
        tbs[0],
        "Dunjun: Procedural Dungeons with In-Run Dynamic Difficulty",
        size=Pt(30),
        bold=True,
        color=WHITE,
    )
    set_shape_text(tbs[1], "Final-Year Project\nYOU SHAOPENG", size=Pt(18), color=MUTED)

    # 2 Overview
    s = prs.slides[1]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Overview", size=Pt(32), bold=True, color=WHITE)
    set_shape_text(
        tbs[1],
        "1. Problem\n"
        "2. The game (PCG + progression)\n"
        "3. Dynamic difficulty design\n"
        "4. Demo\n"
        "5. Results",
        size=Pt(22),
        color=TEXT,
    )

    # 3 Problem
    s = prs.slides[2]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Problem", size=Pt(32), bold=True, color=WHITE)
    set_shape_text(
        tbs[1],
        "Start difficulty does not track in-run struggle",
        size=Pt(16),
        color=MUTED,
    )
    set_shape_text(
        tbs[2],
        "• Start difficulty stays fixed for the whole run\n"
        "• In-run struggle ignored (HP / heals / early damage)\n"
        "• Fixed drops & spikes → harsh or luck-heavy\n"
        "• Hidden reward boosts hard to explain or test",
        size=Pt(20),
        color=TEXT,
    )

    # 4 Goal
    s = prs.slides[3]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Goal", size=Pt(32), bold=True, color=WHITE)
    set_shape_text(tbs[1], "What I built", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[2],
        "In-run dynamic difficulty — explainable rules\n"
        "\n"
        "• Restrained intervention + cooldown\n"
        "• No loot jackpot\n"
        "• ON / OFF for comparison",
        size=Pt(18),
        color=TEXT,
    )

    insert_game_slide(prs, after_index=3)

    # 6 Method (index 5)
    s = prs.slides[5]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Method — 3 Parts", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(tbs[1], "1. Starting tier", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[2],
        "Easy / Medium / Hard\n\nControls how much aid is allowed\n\nEasy → more\nHard → almost none",
        size=Pt(15),
        color=TEXT,
    )
    set_shape_text(tbs[3], "2. Distress value", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[4],
        "Not only “1 HP”\n\nMulti-signal → 0–1\n\nDecides when to act",
        size=Pt(15),
        color=TEXT,
    )
    set_shape_text(tbs[5], "3. Game changes", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[6],
        "When high + not on cooldown:\n\n• potions ↑\n• spikes milder\n• lethal → 1 HP",
        size=Pt(15),
        color=TEXT,
    )

    # 7 Signals
    s = prs.slides[6]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Distress Signals", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(
        tbs[1], "Several signals → one distress value (0–1)", size=Pt(16), color=MUTED
    )
    pairs = [
        ("Health ratio", "Current HP / max HP"),
        ("Room damage", "Damage in room / depth"),
        ("No-heal streak", "Rooms without healing"),
        ("Near-death events", "Near-death count"),
    ]
    idx = 2
    for title, body in pairs:
        set_shape_text(tbs[idx], title, size=Pt(18), bold=True, color=ACCENT)
        set_shape_text(tbs[idx + 1], body, size=Pt(16), color=TEXT)
        idx += 2

    # 8 What Changes
    s = prs.slides[7]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "What Changes", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(tbs[1], "Adjusts", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[2],
        "• Potion chance ↑\n• Spike pressure ↓\n• Easier layouts\n• Lethal → 1 HP (rare)\n• Cooldown",
        size=Pt(17),
        color=TEXT,
    )
    set_shape_text(tbs[3], "Does not", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[4],
        "• Heavy gold boost\n• Relic jackpot\n• Flip the whole run",
        size=Pt(17),
        color=TEXT,
    )

    # 9 Switch
    s = prs.slides[8]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Playtest Switch", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(tbs[1], "Same game, one comparison", size=Pt(16), color=MUTED)
    set_shape_text(tbs[2], "ON", size=Pt(20), bold=True, color=ACCENT)
    set_shape_text(tbs[3], "System active\n→ experiment", size=Pt(18), color=TEXT)
    set_shape_text(tbs[4], "OFF", size=Pt(20), bold=True, color=ACCENT)
    set_shape_text(tbs[5], "System off\n→ control", size=Pt(18), color=TEXT)

    # 10 Demo
    s = prs.slides[9]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Demo", size=Pt(40), bold=True, color=WHITE)
    set_shape_text(tbs[1], "ON (Easy) → short OFF control", size=Pt(20), color=MUTED)

    # 11 Results
    s = prs.slides[10]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Results", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(tbs[1], "Small playtest — exploratory", size=Pt(16), color=MUTED)
    set_shape_text(tbs[2], "Works", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[3],
        "• Integrated + toggle\n• More survivability aid\n• Milder spikes\n• Demo-ready prototype",
        size=Pt(17),
        color=TEXT,
    )
    set_shape_text(tbs[4], "Limits", size=Pt(18), bold=True, color=ACCENT)
    set_shape_text(
        tbs[5],
        "• Small sample\n• Modest gain\n• Light tuning only\n• No HP-only ablation\n• High run variance",
        size=Pt(17),
        color=TEXT,
    )

    # 12 Takeaways
    s = prs.slides[11]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Takeaways", size=Pt(28), bold=True, color=WHITE)
    set_shape_text(
        tbs[1],
        "• Start difficulty ≠ in-run struggle\n"
        "• PCG dungeon + restrained dynamic difficulty\n"
        "• End-to-end, modest gains for now\n"
        "• Next: larger playtest, finer tuning",
        size=Pt(20),
        color=TEXT,
    )

    # 13 Thanks
    s = prs.slides[12]
    tbs = textboxes(s)
    set_shape_text(tbs[0], "Thank you", size=Pt(40), bold=True, color=WHITE)

    alt = FOLDER / "Dynamic_Difficulty_Video_Slides_AUTOSAVE_updated.pptx"
    try:
        prs.save(str(path))
        print(f"Saved into autosave file: {path.name}")
    except PermissionError:
        prs.save(str(alt))
        print(f"Autosave is open/locked. Saved helper file:\n  {alt}")
        # Try overwrite again via copy
        try:
            shutil.copy2(alt, path)
            print(f"Copied over locked path successfully: {path.name}")
        except Exception as e:
            print(f"Could not overwrite autosave yet ({e}).")
            print("Please CLOSE that PowerPoint tab, then tell me — I will write into it.")


if __name__ == "__main__":
    main()
