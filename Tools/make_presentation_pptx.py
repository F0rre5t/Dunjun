"""Generate plain, condensed slides for the final-year project video.

Slides = cues for the viewer. Narration carries the detail (not on-screen subtitles).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "Presentation" / "Dynamic_Difficulty_Video_Slides.pptx"

BG = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SOFT = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xCC, 0xCC, 0xCC)
FONT = "Calibri"


def set_run(run, size=20, bold=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG


def add_rule(slide, top=1.15):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Pt(1)
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = LINE


def add_title(slide, text, top=0.45, size=28):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(0.6))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, size=size, bold=True, color=INK)


def add_lines(slide, lines, left=0.8, top=1.55, width=11.7, height=5.2, size=22, gap=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=INK)


def add_two_cols(slide, left_title, left_lines, right_title, right_lines, top=1.55):
    lt = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(5.5), Inches(0.4))
    r = lt.text_frame.paragraphs[0].add_run()
    r.text = left_title
    set_run(r, size=20, bold=True)

    add_lines(slide, left_lines, left=0.8, top=top + 0.55, width=5.5, height=4.5, size=20, gap=12)

    rt = slide.shapes.add_textbox(Inches(7.0), Inches(top), Inches(5.5), Inches(0.4))
    r = rt.text_frame.paragraphs[0].add_run()
    r.text = right_title
    set_run(r, size=20, bold=True)

    add_lines(slide, right_lines, left=7.0, top=top + 0.55, width=5.5, height=4.5, size=20, gap=12)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Dunjun: Procedural Dungeons with In-Run Dynamic Difficulty"
    set_run(r, size=28, bold=True)

    box2 = s.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.8))
    tf = box2.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Final-Year Project"
    set_run(r, size=18, color=SOFT)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "[Your Name]"
    set_run(r, size=16, color=SOFT)

    # 2 Overview
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Overview")
    add_rule(s)
    add_lines(
        s,
        [
            "1. Problem",
            "2. The game (PCG + progression)",
            "3. Dynamic difficulty design",
            "4. Demo",
            "5. Results",
        ],
        size=24,
        gap=16,
    )

    # 3 Problem
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Problem")
    add_rule(s)
    add_lines(
        s,
        [
            "• Start difficulty stays fixed for the whole run",
            "• In-run struggle is ignored (HP / heals / early damage)",
            "• Fixed drops & spikes → harsh or luck-heavy",
            "• Hidden reward boosts are hard to explain or test",
        ],
        size=24,
        gap=18,
    )

    # 4 Goal
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Goal")
    add_rule(s)
    add_lines(
        s,
        [
            "In-run dynamic difficulty — explainable rules",
            "",
            "• Restrained intervention + cooldown",
            "• No loot jackpot",
            "• ON / OFF for comparison",
        ],
        size=24,
        gap=16,
    )

    # 5 The game (PCG + progression, one slide)
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "The Game")
    add_rule(s)
    add_lines(
        s,
        [
            "PCG",
            "• Random-walk room layout · walls / doors from neighbours",
            "• Start / shop / boss by rules · enemies by depth",
            "• Spike chance & shape scale with progress",
            "",
            "Progression",
            "• Relics + gold shops → different builds each run",
            "• Route, build, and resource pace all vary → state swings hard",
        ],
        size=20,
        gap=8,
    )

    # 6 Method — 3 parts
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Method — 3 Parts")
    add_rule(s)
    add_lines(
        s,
        [
            "1. Starting tier → how much aid is allowed",
            "2. Distress (0–1) → when to intervene",
            "3. Game changes → potions / spikes / lethal save",
        ],
        size=24,
        gap=22,
    )

    # 7 Signals
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Distress Signals")
    add_rule(s)
    add_lines(
        s,
        [
            "• Health ratio",
            "• Room damage",
            "• No-heal streak",
            "• Near-death events",
        ],
        size=26,
        gap=18,
    )

    # 8 What changes
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "What Changes")
    add_rule(s)
    add_two_cols(
        s,
        "Adjusts",
        [
            "• Potion chance ↑",
            "• Spike pressure ↓",
            "• Easier layouts",
            "• Lethal → 1 HP (rare)",
            "• Cooldown",
        ],
        "Does not",
        [
            "• Heavy gold boost",
            "• Relic jackpot",
            "• Flip the whole run",
        ],
        top=1.5,
    )

    # 9 ON/OFF
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Playtest Switch")
    add_rule(s)
    add_two_cols(
        s,
        "ON",
        [
            "System active",
            "→ experiment",
        ],
        "OFF",
        [
            "System off",
            "→ control",
        ],
        top=2.2,
    )

    # 10 Demo
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.8))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "Demo"
    set_run(r, size=40, bold=True)

    box2 = s.shapes.add_textbox(Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.8))
    r = box2.text_frame.paragraphs[0].add_run()
    r.text = "ON (Easy) → short OFF control"
    set_run(r, size=20, color=SOFT)

    # 11 Results
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Results")
    add_rule(s)
    add_two_cols(
        s,
        "Works",
        [
            "• Integrated + toggle",
            "• More survivability aid",
            "• Milder spikes",
            "• Demo-ready prototype",
        ],
        "Limits",
        [
            "• Small sample",
            "• Modest gain",
            "• Light tuning only",
            "• No HP-only ablation",
            "• High run variance",
        ],
        top=1.5,
    )

    # 12 Takeaways
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_title(s, "Takeaways")
    add_rule(s)
    add_lines(
        s,
        [
            "• Start difficulty ≠ in-run struggle",
            "• PCG dungeon + restrained dynamic difficulty",
            "• End-to-end, modest gains for now",
            "• Next: larger playtest, finer tuning",
        ],
        size=24,
        gap=18,
    )

    # 13 Thanks
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    box = s.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.8))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "Thank you"
    set_run(r, size=40, bold=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUT)
        print(f"Saved: {OUT}")
    except PermissionError:
        alt = OUT.with_name(OUT.stem + "_new.pptx")
        prs.save(alt)
        print(f"Original locked. Saved: {alt}")
        print("Close PowerPoint, then re-run to overwrite the main file.")


if __name__ == "__main__":
    main()
